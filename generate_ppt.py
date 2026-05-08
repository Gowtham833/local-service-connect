import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def set_font(run, size, color, bold=False):
    run.font.name = 'Poppins'
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold

prs = Presentation()
# Set slide dimensions to widescreen (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_COLOR = RGBColor(0xFA, 0xF4, 0xDF)
PRIMARY_COLOR = RGBColor(0xE3, 0x56, 0x2B)
SECONDARY_COLOR = RGBColor(0x1D, 0x36, 0x39)
NEUTRAL_COLOR = RGBColor(0x7F, 0x7F, 0x7F)

# Helper to set background
def set_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

# Slide 1: Cover
slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
set_bg(slide)

txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(1.5))
tf = txBox.text_frame
p = tf.add_paragraph()
run = p.add_run()
run.text = "Local Service Connect"
set_font(run, 44, SECONDARY_COLOR, True)

p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = "Serverless AI-Driven Service Marketplace"
set_font(run2, 24, NEUTRAL_COLOR)

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(3))
tf2 = txBox2.text_frame
p3 = tf2.add_paragraph()
run3 = p3.add_run()
run3.text = "Team Name: Team Local Service Connect"
set_font(run3, 22, PRIMARY_COLOR, True)

p4 = tf2.add_paragraph()
run4 = p4.add_run()
run4.text = "Team Members:"
set_font(run4, 20, SECONDARY_COLOR, True)

members = [
    "Bhavani Gandhavalla (Team Lead)",
    "Gorle Rahul",
    "Mohan Sivaram",
    "Kota Gowtham (Infrastructure Lead)",
    "Yeruva Tarun",
    "Jajimoggala Sai"
]

for m in members:
    p_m = tf2.add_paragraph()
    p_m.level = 1
    run_m = p_m.add_run()
    run_m.text = m
    set_font(run_m, 18, SECONDARY_COLOR)


# Function to create standard content slides
def create_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = title
    set_font(run, 36, SECONDARY_COLOR, True)
    
    # Content
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
    tf2 = txBox2.text_frame
    
    for bullet in bullets:
        p_b = tf2.add_paragraph()
        p_b.level = 0
        run_b = p_b.add_run()
        run_b.text = bullet
        set_font(run_b, 20, SECONDARY_COLOR)

# Slide 2: Problem & Solution Overview
create_content_slide("Problem & Solution Overview", [
    "Problem: Finding reliable local services with transparent pricing is difficult and time-consuming.",
    "Problem: There is high manual effort required to connect customers with available service workers.",
    "Solution: Local Service Connect is a serverless cloud platform connecting customers to workers.",
    "Solution: Integrated GenAI automates matching, price estimation, and customer support."
])

# Slide 3: Core Features & Workflow
create_content_slide("Core Features & Workflow", [
    "Features: AI-powered worker matching and real-time price estimation (INR).",
    "Features: Automated customer support chatbot (ServiBot) and review sentiment analysis.",
    "Workflow: Customer requests a service via the secure web portal.",
    "Workflow: AI estimates the cost and matches the best available worker.",
    "Workflow: Worker accepts, completes the job, and the customer provides feedback."
])

# Slide 4: Technology Stack
create_content_slide("Technology Stack & Architecture", [
    "Frontend: HTML, CSS, JavaScript (Hosted on Amazon S3 & CloudFront).",
    "Backend: Node.js API running on AWS ECS Fargate (Serverless compute).",
    "Database: Relational data management using Amazon RDS (PostgreSQL).",
    "Cloud & AI: VPC, Application Load Balancer, Amazon Bedrock (Claude 3), Comprehend.",
    "Infrastructure: Infrastructure as Code (IaC) using Terraform."
])

# Slide 5: Team Learnings & Conclusion
create_content_slide("Team Learnings & Conclusion", [
    "Team Intro: A 6-member team specializing across Cloud Infra, Backend, Frontend, AI/ML, DB, and Security.",
    "Learnings: Successfully designed and deployed a highly available serverless AWS architecture.",
    "Learnings: Integrated advanced Generative AI models to solve real-world marketplace challenges.",
    "Learnings: Collaborated effectively using Git workflows and automated Terraform deployments.",
    "Thank You! Questions and Feedback?"
])

prs.save("Local_Service_Connect_Presentation.pptx")
print("Presentation saved to Local_Service_Connect_Presentation.pptx")
